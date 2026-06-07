import os
import shutil
from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form
from fastapi.responses import FileResponse, HTMLResponse
from sqlalchemy.orm import Session
from typing import List, Optional
from app.database import get_db
from app import models
from app.schemas.document import DocumentCreate, DocumentUpdate, DocumentOut
from app.utils.auth import get_current_user
from app.config import UPLOAD_DIR

router = APIRouter(prefix="/documents", tags=["documents"])


@router.get("/", response_model=List[DocumentOut])
def list_documents(db: Session = Depends(get_db), current_user: models.User = Depends(get_current_user)):
    return db.query(models.Document).filter(models.Document.user_id == current_user.id).all()


@router.post("/", response_model=DocumentOut)
async def create_document(
    title: str = Form(...),
    doc_type: str = Form(...),
    notion_url: Optional[str] = Form(None),
    file: Optional[UploadFile] = File(None),
    db: Session = Depends(get_db),
    current_user: models.User = Depends(get_current_user)
):
    doc = models.Document(
        user_id=current_user.id,
        title=title,
        doc_type=doc_type,
        notion_url=notion_url,
    )

    if file and file.filename:
        doc_dir = os.path.join(UPLOAD_DIR, "documents", str(current_user.id))
        os.makedirs(doc_dir, exist_ok=True)
        file_path = os.path.join(doc_dir, file.filename)
        with open(file_path, "wb") as f:
            shutil.copyfileobj(file.file, f)
        file_size = os.path.getsize(file_path)
        doc.file_path = file_path
        doc.file_name = file.filename
        doc.file_size = file_size

    db.add(doc)
    db.commit()
    db.refresh(doc)
    return doc


@router.put("/{doc_id}", response_model=DocumentOut)
def update_document(doc_id: int, data: DocumentUpdate, db: Session = Depends(get_db),
                    current_user: models.User = Depends(get_current_user)):
    doc = db.query(models.Document).filter(
        models.Document.id == doc_id, models.Document.user_id == current_user.id
    ).first()
    if not doc:
        raise HTTPException(status_code=404, detail="문서를 찾을 수 없습니다.")

    if data.title is not None:
        doc.title = data.title
    if data.doc_type is not None:
        doc.doc_type = data.doc_type
    if data.notion_url is not None:
        doc.notion_url = data.notion_url

    db.commit()
    db.refresh(doc)
    return doc


@router.delete("/{doc_id}")
def delete_document(doc_id: int, db: Session = Depends(get_db),
                    current_user: models.User = Depends(get_current_user)):
    doc = db.query(models.Document).filter(
        models.Document.id == doc_id, models.Document.user_id == current_user.id
    ).first()
    if not doc:
        raise HTTPException(status_code=404, detail="문서를 찾을 수 없습니다.")

    if doc.file_path and os.path.exists(doc.file_path):
        os.remove(doc.file_path)

    db.delete(doc)
    db.commit()
    return {"message": "삭제되었습니다."}


@router.get("/{doc_id}/download")
def download_document(doc_id: int, db: Session = Depends(get_db),
                      current_user: models.User = Depends(get_current_user)):
    doc = db.query(models.Document).filter(
        models.Document.id == doc_id, models.Document.user_id == current_user.id
    ).first()
    if not doc or not doc.file_path:
        raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")
    if not os.path.exists(doc.file_path):
        raise HTTPException(status_code=404, detail="파일이 서버에 존재하지 않습니다.")

    return FileResponse(
        path=doc.file_path,
        filename=doc.file_name,
        media_type="application/octet-stream"
    )


@router.get("/{doc_id}/view")
def view_document(doc_id: int, db: Session = Depends(get_db),
                  current_user: models.User = Depends(get_current_user)):
    doc = db.query(models.Document).filter(
        models.Document.id == doc_id, models.Document.user_id == current_user.id
    ).first()
    if not doc or not doc.file_path:
        raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")
    if not os.path.exists(doc.file_path):
        raise HTTPException(status_code=404, detail="파일이 서버에 존재하지 않습니다.")

    ext = doc.file_name.rsplit(".", 1)[-1].lower() if doc.file_name else ""
    media_type_map = {
        "pdf": "application/pdf",
        "jpg": "image/jpeg", "jpeg": "image/jpeg",
        "png": "image/png", "gif": "image/gif",
        "txt": "text/plain",
    }
    media_type = media_type_map.get(ext, "application/octet-stream")

    return FileResponse(path=doc.file_path, media_type=media_type)


@router.get("/{doc_id}/html-preview")
def html_preview_document(doc_id: int, db: Session = Depends(get_db),
                           current_user: models.User = Depends(get_current_user)):
    doc = db.query(models.Document).filter(
        models.Document.id == doc_id, models.Document.user_id == current_user.id
    ).first()
    if not doc or not doc.file_path or not os.path.exists(doc.file_path):
        raise HTTPException(status_code=404, detail="파일을 찾을 수 없습니다.")

    ext = (doc.file_name.rsplit(".", 1)[-1].lower()) if doc.file_name else ""

    if ext in ("docx", "doc"):
        try:
            from docx import Document as DocxDocument
            from docx.shared import Pt
            import html as html_mod
            document = DocxDocument(doc.file_path)
            body = ""
            for para in document.paragraphs:
                if para.text.strip():
                    style = para.style.name if para.style else ""
                    text = html_mod.escape(para.text)
                    if "Heading 1" in style:
                        body += f"<h1>{text}</h1>"
                    elif "Heading 2" in style:
                        body += f"<h2>{text}</h2>"
                    elif "Heading 3" in style:
                        body += f"<h3>{text}</h3>"
                    else:
                        body += f"<p>{text}</p>"
                else:
                    body += "<br>"
            content = f"""<!DOCTYPE html><html><head><meta charset="utf-8">
<style>body{{font-family:sans-serif;padding:24px;line-height:1.7;color:#333;max-width:800px;margin:0 auto}}
h1,h2,h3{{color:#111}}p{{margin:.5em 0}}</style></head><body>{body}</body></html>"""
            return HTMLResponse(content=content)
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"DOCX 변환 실패: {str(e)}")

    elif ext in ("xlsx", "xls"):
        try:
            import openpyxl
            import html as html_mod
            wb = openpyxl.load_workbook(doc.file_path, read_only=True, data_only=True)
            tabs = ""
            sheets_html = ""
            for i, ws in enumerate(wb.worksheets):
                tab_id = f"sheet{i}"
                tabs += f'<button onclick="showSheet(\'{tab_id}\')" id="btn_{tab_id}" class="tab">{html_mod.escape(ws.title)}</button>'
                rows_html = ""
                for row in ws.iter_rows(values_only=True):
                    cells = "".join(f"<td>{html_mod.escape(str(c)) if c is not None else ''}</td>" for c in row)
                    rows_html += f"<tr>{cells}</tr>"
                display = "block" if i == 0 else "none"
                sheets_html += f'<div id="{tab_id}" class="sheet" style="display:{display}"><table>{rows_html}</table></div>'
            content = f"""<!DOCTYPE html><html><head><meta charset="utf-8">
<style>body{{font-family:sans-serif;padding:16px;font-size:13px}}
.tab{{padding:6px 14px;margin-right:4px;cursor:pointer;border:1px solid #ccc;border-radius:6px 6px 0 0;background:#f5f5f5}}
.tab.active{{background:#4F46E5;color:#fff;border-color:#4F46E5}}
table{{border-collapse:collapse;width:100%}}
td{{border:1px solid #ddd;padding:4px 8px;white-space:nowrap}}
tr:nth-child(1){{background:#f0f0f0;font-weight:bold}}</style></head>
<body><div style="margin-bottom:8px">{tabs}</div>{sheets_html}
<script>function showSheet(id){{document.querySelectorAll('.sheet').forEach(s=>s.style.display='none');
document.getElementById(id).style.display='block';
document.querySelectorAll('.tab').forEach(b=>b.classList.remove('active'));
document.getElementById('btn_'+id).classList.add('active')}}
document.querySelector('.tab')&&document.querySelector('.tab').classList.add('active')</script>
</body></html>"""
            return HTMLResponse(content=content)
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"XLSX 변환 실패: {str(e)}")

    elif ext in ("pptx", "ppt"):
        try:
            from pptx import Presentation
            from pptx.util import Inches
            import html as html_mod
            prs = Presentation(doc.file_path)
            slides_html = ""
            for i, slide in enumerate(prs.slides, 1):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = " ".join(run.text for run in para.runs if run.text.strip())
                            if line:
                                texts.append(html_mod.escape(line))
                slide_content = "".join(f"<p>{t}</p>" for t in texts) or "<p class='empty'>내용 없음</p>"
                slides_html += f'<div class="slide"><div class="slide-num">{i}슬라이드</div>{slide_content}</div>'
            content = f"""<!DOCTYPE html><html><head><meta charset="utf-8">
<style>body{{font-family:sans-serif;padding:16px;background:#f5f5f5}}
.slide{{background:#fff;border-radius:8px;padding:20px 24px;margin-bottom:12px;box-shadow:0 1px 4px rgba(0,0,0,.1)}}
.slide-num{{font-size:11px;color:#999;margin-bottom:8px;font-weight:600}}
p{{margin:.4em 0;font-size:14px;line-height:1.6;color:#333}}
.empty{{color:#bbb;font-style:italic}}</style></head>
<body>{slides_html}</body></html>"""
            return HTMLResponse(content=content)
        except Exception as e:
            raise HTTPException(status_code=500, detail=f"PPT 변환 실패: {str(e)}")

    else:
        raise HTTPException(status_code=400, detail=f"미리보기 미지원: {ext.upper()}")
